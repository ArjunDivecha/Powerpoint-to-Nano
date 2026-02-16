#!/usr/bin/env python3

"""streamlit_app.py

Purpose
- Local Streamlit web UI for the PowerPoint-to-Nano pipeline.
- Lets you:
  - Pick various input files using a native macOS file picker (AppleScript/osascript)
  - Render slides/pages via LibreOffice (PPTX) or conversion utilities (other formats)
  - Generate / regenerate one slide at a time in a chosen style (Interactions API)
  - Auto-detect text-heavy slides for enhanced text extraction accuracy
  - Preview a PDF in-browser before saving
  - Save the PDF next to the input file as: <stem>_nano.pdf (overwrite)

INPUT FILES (prominent)
- Multiple formats supported via the in-app picker:
  - PPTX (PowerPoint presentations)
  - PDF (multi-page documents)
  - GIF (animated graphics)
  - TXT (plain text files)
  - DOCX (Word documents)
  - Images (PNG, JPG, JPEG, WebP, etc.)
  - Markdown (.md files)

OUTPUT FILES (prominent)
- Cached style example images (generated once per style):
  - {repo}/style_examples_cache/slide1/<style>.png
  - {repo}/style_examples_cache/slide2/<style>.png
- Rendered slide images (LibreOffice exports or conversions):
  - {repo}/pptx2nano_output_streamlit/{deck_stem}/rendered/*.png
- Final PDF (written only when you click Save):
  - {input_folder}/{input_stem}_nano.pdf

Version History
- v0.1.0 (2025-12-13): Initial Streamlit UI with PPTX support only.
- v0.2.0 (2026-01-03): Added multi-format support, text extraction, auto-detection, custom styles
- v0.3.0 (2026-02-16): Streamlit PPTX path standardized on LibreOffice; docs aligned with current behavior

Last Updated
- 2026-02-16

Features
- 33+ built-in styles including ghibli-pride
- Custom style support with on-the-fly generation
- Text extraction for PPTX files with auto-detection
- Text deduplication to prevent duplicate titles
- Fresh regenerations (always from original, not previous generation)
- Parallel "Generate ALL" with configurable worker count
- Faster preview rendering with configurable max pages

Platform Note
- The "Choose file..." button currently uses AppleScript (`osascript`), which is macOS-specific.
"""

from __future__ import annotations

import base64
import concurrent.futures
import re
import subprocess
import tempfile
import time
from html import unescape
from io import BytesIO
from pathlib import Path

import streamlit as st
from PIL import Image
from pptx import Presentation

import pptx2nano
from pptx_converter import export_slides


def _pick_input_file() -> Path | None:
    # IMPORTANT: On macOS, Streamlit runs this script in a worker thread.
    # tkinter uses AppKit and will crash with:
    #   "NSWindow should only be instantiated on the main thread!"
    # Instead we use AppleScript's "choose file" via osascript.

    # Return POSIX path directly, and handle user cancel (AppleScript error -128)
    # by returning an empty string.
    applescript = r'''try
  set f to choose file with prompt "Select a file (.pptx, .pdf, .gif, .txt, .docx, .png, .jpg, .jpeg, .webp, .md)"
  return POSIX path of f
on error number -128
  return ""
end try'''

    proc = subprocess.run(
        ["osascript", "-e", applescript],
        check=True,
        capture_output=True,
        text=True,
    )

    selected = (proc.stdout or "").strip()
    if not selected:
        return None
    return Path(selected).expanduser().resolve()


def _ensure_session_defaults() -> None:
    st.session_state.setdefault("input_path", None)
    st.session_state.setdefault("input_type", None)
    st.session_state.setdefault("rendered_paths", None)
    st.session_state.setdefault("generated_images", {})
    st.session_state.setdefault("interaction_ids", {})
    st.session_state.setdefault("last_preview_pdf", None)
    st.session_state.setdefault("style_example_set", "slide1")
    st.session_state.setdefault("selected_style", "lego")
    st.session_state.setdefault("pptx_conversion_method", "libreoffice")  # Default to LibreOffice
    st.session_state.setdefault("selected_slide", 1)
    st.session_state.setdefault("pptx_text_cache", {})
    st.session_state.setdefault("text_extraction_mode", "strict")
    st.session_state.setdefault("dedupe_extracted_text", False)
    st.session_state.setdefault("generate_workers", 4)
    st.session_state.setdefault("preview_max_pages", 10)
    st.session_state.setdefault("libreoffice_timeout_seconds", 120)


def _reset_for_new_input(input_path: Path, input_type: str) -> None:
    st.session_state["input_path"] = input_path
    st.session_state["input_type"] = input_type
    st.session_state["rendered_paths"] = None
    st.session_state["generated_images"] = {}
    st.session_state["interaction_ids"] = {}
    st.session_state["last_preview_pdf"] = None
    st.session_state["pptx_text_cache"] = {}


def _pptx_cache_key(pptx_path: Path) -> str:
    stat = pptx_path.stat()
    return f"{pptx_path.resolve()}:{stat.st_size}:{stat.st_mtime_ns}"


def _extract_all_text_from_pptx(pptx_path: Path) -> list[str]:
    prs = Presentation(pptx_path)
    texts: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape_text = shape.text.strip()
                if shape_text:
                    parts.append(shape_text)
        texts.append("\n".join(parts))
    return texts


def _get_pptx_text_cache(pptx_path: Path) -> list[str]:
    cache = st.session_state.setdefault("pptx_text_cache", {})
    key = _pptx_cache_key(pptx_path)
    if cache.get("key") != key:
        cache = {"key": key, "slides": _extract_all_text_from_pptx(pptx_path)}
        st.session_state["pptx_text_cache"] = cache
    return list(cache.get("slides", []))


def _dedupe_text_lines(text: str) -> str:
    seen: set[str] = set()
    unique_lines: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped in seen:
            continue
        seen.add(stripped)
        unique_lines.append(line)
    return "\n".join(unique_lines)


def _wrap_line_to_width(draw, text: str, font, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]

    lines: list[str] = []
    current: list[str] = []

    for word in words:
        candidate = " ".join(current + [word])
        bbox = draw.textbbox((0, 0), candidate, font=font)
        if (bbox[2] - bbox[0]) <= max_width or not current:
            current.append(word)
        else:
            lines.append(" ".join(current))
            current = [word]

    if current:
        lines.append(" ".join(current))
    return lines


def _render_text_tokens_to_pages(
    tokens: list[tuple[str, bool]],
    output_dir: Path,
    prefix: str,
) -> list[Path]:
    """Render (text, is_header) tokens into one or more 1920x1080 PNG pages."""
    from PIL import ImageDraw, ImageFont

    output_dir.mkdir(parents=True, exist_ok=True)

    img_width, img_height = 1920, 1080
    margin = 100
    max_width = img_width - 2 * margin
    page_bottom = img_height - margin

    try:
        body_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
        header_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 50)
    except Exception:
        body_font = ImageFont.load_default()
        header_font = body_font

    line_step_body = 60
    line_step_header = 70

    def new_page():
        page = Image.new("RGB", (img_width, img_height), color="white")
        return page, ImageDraw.Draw(page), margin

    rendered_paths: list[Path] = []
    page_num = 1
    page_img, draw, y_pos = new_page()

    def flush_page(img_obj, path_num: int) -> None:
        out_path = output_dir / f"{prefix}_{path_num:03d}.png"
        img_obj.save(out_path, "PNG")
        rendered_paths.append(out_path)

    has_drawn_any_text = False
    for text, is_header in tokens:
        if text == "":
            y_pos += line_step_body // 2
            continue

        font = header_font if is_header else body_font
        step = line_step_header if is_header else line_step_body
        for wrapped in _wrap_line_to_width(draw, text, font, max_width):
            if y_pos + step > page_bottom:
                flush_page(page_img, page_num)
                page_num += 1
                page_img, draw, y_pos = new_page()
            draw.text((margin, y_pos), wrapped, fill="black", font=font)
            y_pos += step
            has_drawn_any_text = True

        if is_header:
            y_pos += 8

    if not has_drawn_any_text:
        draw.text((margin, margin), "(Empty input)", fill="black", font=body_font)
    flush_page(page_img, page_num)
    return rendered_paths


def _markdown_to_tokens(md_content: str) -> list[tuple[str, bool]]:
    """Convert markdown to display tokens with basic structural fidelity."""
    import markdown

    html = markdown.markdown(md_content, extensions=["extra", "fenced_code", "tables", "nl2br"])

    # Preserve structural boundaries before stripping tags.
    html = re.sub(r"<h([1-6])>", r"\n__P2N_H\1__", html, flags=re.IGNORECASE)
    html = re.sub(r"</h[1-6]>", "\n", html, flags=re.IGNORECASE)
    html = re.sub(r"<li>", "\n__P2N_LI__", html, flags=re.IGNORECASE)
    html = re.sub(r"</li>", "\n", html, flags=re.IGNORECASE)
    html = re.sub(r"<br\\s*/?>", "\n", html, flags=re.IGNORECASE)
    html = re.sub(r"</(p|div|ul|ol|pre|table|tr|blockquote)>", "\n", html, flags=re.IGNORECASE)
    html = re.sub(r"<[^>]+>", "", html)
    text = unescape(html)

    tokens: list[tuple[str, bool]] = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            if tokens and tokens[-1][0] != "":
                tokens.append(("", False))
            continue
        if line.startswith("__P2N_H"):
            line = re.sub(r"^__P2N_H[1-6]__", "", line).strip()
            if line:
                tokens.append((line, True))
            continue
        if line.startswith("__P2N_LI__"):
            line = line.replace("__P2N_LI__", "", 1).strip()
            tokens.append((f"• {line}" if line else "•", False))
            continue
        tokens.append((line, False))

    if not tokens:
        tokens = [("(Empty markdown)", False)]
    return tokens


def _style_options() -> list[str]:
    return sorted(list(pptx2nano.BUILTIN_STYLES.keys())) + ["custom"]


def _get_selected_style(style_choice: str, custom_style: str) -> str | None:
    if style_choice == "custom":
        s = (custom_style or "").strip()
        return s if s else None
    return style_choice


def _style_example_cache_root() -> Path:
    return Path(__file__).resolve().parent / "style_examples_cache"


def _style_example_cache_base(style_example_set: str) -> Path:
    root = _style_example_cache_root()
    if style_example_set == "slide2":
        return root / "slide2"
    if style_example_set == "slide1":
        return root / "slide1"
    return root


def _style_example_cache_path(style: str, *, style_example_set: str) -> Path:
    safe = "".join(ch for ch in style.lower() if ch.isalnum() or ch in ("-", "_"))
    if not safe:
        safe = "style"
    return _style_example_cache_base(style_example_set) / f"{safe}.png"


def _get_or_create_style_example(style: str, image_model: str, *, style_example_set: str) -> bytes:
    cache_path = _style_example_cache_path(style, style_example_set=style_example_set)
    if cache_path.exists():
        return cache_path.read_bytes()

    # Backward-compatible fallback: if slide1 set is selected but hasn't been copied yet,
    # fall back to the legacy root cache.
    if style_example_set == "slide1":
        safe = "".join(ch for ch in style.lower() if ch.isalnum() or ch in ("-", "_"))
        if not safe:
            safe = "style"
        legacy_path = _style_example_cache_root() / f"{safe}.png"
        if legacy_path.exists():
            return legacy_path.read_bytes()

    # We do not auto-generate missing thumbnails for slide1/slide2 sets for BUILTIN styles,
    # because those sets are meant to be precomputed from specific deck slides.
    # However, for custom styles, we allow on-the-fly generation.
    is_builtin_style = style in pptx2nano.BUILTIN_STYLES
    if style_example_set in {"slide1", "slide2"} and is_builtin_style:
        raise RuntimeError(
            f"Missing cached example for style='{style}' in set='{style_example_set}'. "
            "Generate the cache first (style-sample-generator)."
        )

    # Legacy fallback (root cache): preserve old behavior (generate-on-miss).
    client = pptx2nano.create_client()
    prompt = f"""You are an expert presentation designer.

TASK
Create a single-page example exhibit in the visual style '{style}'.

CONSTRAINTS (VERY IMPORTANT)
- Make it look like a real, clean infographic exhibit.
- Use clear typography and spacing.
- Include a simple chart and a small table with fake placeholder values.
- Do NOT mention this is a sample.

OUTPUT
- Generate exactly ONE image.
""".strip()

    interaction = client.interactions.create(
        model=image_model,
        input=prompt,
        response_modalities=["IMAGE"],
    )

    img_bytes, _mime = pptx2nano._extract_image_bytes_from_interaction(interaction)
    img = Image.open(BytesIO(img_bytes))
    out = BytesIO()
    img.save(out, format="PNG")
    cache_path.parent.mkdir(parents=True, exist_ok=True)
    cache_path.write_bytes(out.getvalue())
    return cache_path.read_bytes()


def _process_pdf_input(pdf_path: Path, output_dir: Path) -> list[Path]:
    """Extract pages from PDF as images."""
    from pdf2image import convert_from_path
    
    output_dir.mkdir(parents=True, exist_ok=True)
    images = convert_from_path(pdf_path, dpi=150)
    
    rendered_paths = []
    for i, img in enumerate(images, 1):
        out_path = output_dir / f"page_{i:03d}.png"
        img.save(out_path, "PNG")
        rendered_paths.append(out_path)
    
    return rendered_paths


def _process_gif_input(gif_path: Path, output_dir: Path) -> list[Path]:
    """Extract frames from GIF as images."""
    output_dir.mkdir(parents=True, exist_ok=True)
    
    with Image.open(gif_path) as img:
        rendered_paths = []
        frame_num = 0
        
        try:
            while True:
                img.seek(frame_num)
                out_path = output_dir / f"frame_{frame_num + 1:03d}.png"
                # Convert to RGB if necessary (GIFs can be in palette mode)
                frame = img.convert("RGB")
                frame.save(out_path, "PNG")
                rendered_paths.append(out_path)
                frame_num += 1
        except EOFError:
            pass
    
    return rendered_paths


def _process_text_input(text_path: Path, output_dir: Path) -> list[Path]:
    """Convert text file to paginated images without truncating overflow."""
    text_content = text_path.read_text(encoding="utf-8")

    paragraphs = [p.strip() for p in text_content.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [p.strip() for p in text_content.splitlines() if p.strip()]
    if not paragraphs:
        paragraphs = ["(Empty text file)"]

    tokens: list[tuple[str, bool]] = []
    for para in paragraphs:
        tokens.append((para, False))
        tokens.append(("", False))
    return _render_text_tokens_to_pages(tokens, output_dir, prefix="text")


def _process_docx_input(docx_path: Path, output_dir: Path) -> list[Path]:
    """Extract text from DOCX and convert to paginated images."""
    from docx import Document

    doc = Document(docx_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        paragraphs = ["(Empty document)"]

    tokens: list[tuple[str, bool]] = []
    for para in paragraphs:
        tokens.append((para, False))
        tokens.append(("", False))
    return _render_text_tokens_to_pages(tokens, output_dir, prefix="docx")


def _process_image_input(image_path: Path, output_dir: Path) -> list[Path]:
    """Process single image file (PNG, JPG, JPEG, WebP, etc.)."""
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Open and convert image to RGB if necessary
    with Image.open(image_path) as img:
        # Convert to RGB if needed (handles RGBA, palette mode, etc.)
        if img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')
        
        # Save as PNG
        out_path = output_dir / "image_001.png"
        img.save(out_path, "PNG")
    
    return [out_path]


def _process_markdown_input(md_path: Path, output_dir: Path) -> list[Path]:
    """Convert Markdown file to paginated images with basic structure preservation."""
    md_content = md_path.read_text(encoding="utf-8")
    tokens = _markdown_to_tokens(md_content)
    return _render_text_tokens_to_pages(tokens, output_dir, prefix="md")


def _extract_text_from_pptx_slide(pptx_path: Path, slide_index: int) -> str:
    """Extract text from a PPTX slide using a cached one-pass parse."""
    try:
        slide_texts = _get_pptx_text_cache(pptx_path)
        if 0 <= slide_index < len(slide_texts):
            return slide_texts[slide_index]
        return ""
    except Exception:
        return ""


def _clean_text_with_gemini(raw_text: str) -> str:
    """Use gemini-3-flash-preview to clean and structure extracted text."""
    if not raw_text.strip():
        return ""
    
    try:
        client = pptx2nano.create_client()
        response = client.models.generate_content(
            model="gemini-3-flash-preview",
            contents=f"""Clean and structure the following slide text. Fix any formatting issues, 
preserve all content, and maintain the original meaning. Return only the cleaned text:

{raw_text}"""
        )
        return response.text.strip()
    except Exception:
        return raw_text


def _should_use_text_extraction(
    rendered_path: Path,
    pptx_path: Path | None,
    slide_index: int,
    extracted_text: str | None = None,
) -> bool:
    """Auto-detect if a slide needs text extraction for better accuracy.
    
    Criteria:
    - Slide has lots of text (detected from PPTX)
    - Image has small dimensions (likely small fonts)
    """
    if not pptx_path or not pptx_path.exists():
        return False
    
    try:
        # Check 1: Extract text and see if it's text-heavy
        text_for_count = extracted_text
        if text_for_count is None:
            text_for_count = _extract_text_from_pptx_slide(pptx_path, slide_index)
        word_count = len(text_for_count.split())
        
        # Check 2: Check image dimensions (small images = small fonts)
        with Image.open(rendered_path) as im:
            width, height = im.size
            is_small = width < 1200 or height < 900
        
        # Use text extraction if:
        # - More than 20 words (lowered threshold for better detection)
        # - OR small image dimensions (likely small fonts)
        return word_count > 20 or is_small
    except Exception:
        return False


def _call_slide_image_model(
    slide_index_1based: int,
    rendered_path: Path,
    image_model: str,
    style: str | None,
    total_slides: int,
    previous_interaction_id: str | None,
    pptx_path: Path | None = None,
    text_extraction_mode: str = "strict",
    dedupe_extracted_text: bool = False,
    extracted_slide_text: str | None = None,
) -> tuple[bytes, str]:
    client = pptx2nano.create_client()

    with Image.open(rendered_path) as im:
        source_width, source_height = im.size

    # Don't use previous_interaction_id for restyling - always start fresh from original slide
    # This ensures dramatic style changes are properly applied

    # Auto-detect if text extraction is needed (unless explicitly disabled).
    use_text_extraction = False
    if text_extraction_mode != "off":
        use_text_extraction = _should_use_text_extraction(
            rendered_path,
            pptx_path,
            slide_index_1based - 1,
            extracted_text=extracted_slide_text,
        )
    
    base_prompt = pptx2nano.build_image_model_prompt(
        slide_index_1based=slide_index_1based,
        total_slides=total_slides,
        source_width=source_width,
        source_height=source_height,
        style=style,
    )

    prompt = base_prompt
    if use_text_extraction and pptx_path:
        raw_text = extracted_slide_text
        if raw_text is None:
            raw_text = _extract_text_from_pptx_slide(pptx_path, slide_index_1based - 1)

        prepared_text = raw_text
        if text_extraction_mode == "assisted":
            prepared_text = _clean_text_with_gemini(prepared_text)
        if dedupe_extracted_text:
            prepared_text = _dedupe_text_lines(prepared_text)

        if prepared_text.strip():
            style_label = style if style else "clean professional"
            constraints = [
                "- Use the EXTRACTED TEXT above for accurate text content.",
                "- Use the IMAGE for layout, visual elements, charts, and spatial relationships.",
                "- Keep the SAME aspect ratio as the input slide.",
                f"- Target pixel size: {source_width} x {source_height}.",
                f"- Apply the '{style_label}' visual style.",
                "- Do NOT add speaker notes.",
                "- Do NOT invent new facts.",
                "- Do NOT duplicate titles, headings, or text unless duplicates are clearly intentional.",
                "- Improve layout, spacing, and readability.",
            ]
            if text_extraction_mode == "strict":
                constraints.insert(
                    1,
                    "- Treat extracted text as authoritative; preserve wording, numbers, and labels exactly.",
                )
            if dedupe_extracted_text:
                constraints.insert(
                    1,
                    "- If the same text appears multiple times in the extracted content, only show it once.",
                )

            # Build enhanced prompt with extracted text
            prompt = f"""You are an expert presentation designer.

TASK
Redesign the provided slide image as a clean, professional infographic slide.

EXTRACTED TEXT CONTENT (use this for accurate text - the image may have small/blurry text):
{prepared_text}

CONSTRAINTS (VERY IMPORTANT)
{chr(10).join(constraints)}

SLIDE CONTEXT
- Slide {slide_index_1based} of {total_slides}

OUTPUT
- Generate exactly ONE image.
""".strip()

    image_bytes = rendered_path.read_bytes()
    mime_type = pptx2nano.mimetypes.guess_type(rendered_path.name)[0] or "application/octet-stream"

    interaction = client.interactions.create(
        model=image_model,
        input=[
            {
                "type": "image",
                "data": base64.b64encode(image_bytes).decode("utf-8"),
                "mime_type": mime_type,
            },
            {"type": "text", "text": prompt},
        ],
        response_modalities=["IMAGE"],
    )

    img_bytes, _mime = pptx2nano._extract_image_bytes_from_interaction(interaction)
    return img_bytes, interaction.id


def _build_pdf_bytes(slide_numbers: list[int]) -> bytes:
    image_bytes_list: list[bytes] = []
    for n in slide_numbers:
        img_bytes = st.session_state["generated_images"].get(n)
        if img_bytes:
            image_bytes_list.append(img_bytes)

    if not image_bytes_list:
        raise RuntimeError("No generated slides available to build a PDF.")

    with tempfile.TemporaryDirectory(prefix="p2n_preview_pdf_") as td:
        tmp_dir = Path(td)
        image_paths: list[Path] = []
        for i, img_bytes in enumerate(image_bytes_list, start=1):
            image_path = tmp_dir / f"preview_{i:03d}.png"
            image_path.write_bytes(img_bytes)
            image_paths.append(image_path)

        out_pdf = tmp_dir / "preview.pdf"
        pptx2nano.images_to_pdf(image_paths, out_pdf)
        return out_pdf.read_bytes()


def _render_pdf_inline(pdf_bytes: bytes, height: int = 800, max_pages: int = 10) -> None:
    # Convert PDF to images for reliable display in Chrome
    try:
        from pdf2image import convert_from_bytes, pdfinfo_from_bytes

        page_cap = max(1, int(max_pages))
        total_pages = 0
        try:
            info = pdfinfo_from_bytes(pdf_bytes)
            total_pages = int(info.get("Pages", 0))
        except Exception:
            total_pages = 0

        last_page = page_cap if total_pages <= 0 else min(total_pages, page_cap)
        images = convert_from_bytes(
            pdf_bytes,
            dpi=150,
            size=(800, None),
            first_page=1,
            last_page=last_page,
        )

        if images:
            shown = len(images)
            if total_pages > 0:
                st.info(f"PDF Preview (showing {shown} of {total_pages} pages):")
            else:
                st.info(f"PDF Preview (showing {shown} pages):")
            for i, img in enumerate(images, 1):
                st.image(img, caption=f"Page {i}", use_container_width=True)
        else:
            raise Exception("No pages found")
            
    except Exception as e:
        # Ultimate fallback
        st.warning("PDF preview not available in this browser. Use the download button below to view the PDF.")
        st.caption(f"Preview error: {str(e)}")


def main() -> None:
    st.set_page_config(page_title="PowerPoint to Nano", layout="wide")
    _ensure_session_defaults()

    st.markdown(
        """
        <div style="padding: 1rem 0 0.5rem 0;">
          <h1 style="margin: 0;">PowerPoint to Nano</h1>
          <p style="margin: 0.25rem 0 0 0; color: #6b7280;">
            Pick a file (PPTX, PDF, GIF, TXT, DOCX, PNG, JPG, MD), pick a style, generate pages, preview the PDF, then save.
          </p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    left, right = st.columns([1, 1])

    with left:
        st.subheader("1) Choose Input File")
        if st.button("Choose file…", use_container_width=True):
            try:
                picked = _pick_input_file()
                if picked is None:
                    st.info("No file selected.")
                else:
                    suffix = picked.suffix.lower()
                    supported_formats = [".pptx", ".pdf", ".gif", ".txt", ".docx", ".png", ".jpg", ".jpeg", ".webp", ".md"]
                    if suffix not in supported_formats:
                        st.error(f"Unsupported file type. Please select: {', '.join(supported_formats)}")
                    else:
                        # Normalize image extensions
                        if suffix in [".png", ".jpg", ".jpeg", ".webp"]:
                            input_type = "image"
                        else:
                            input_type = suffix[1:]  # Remove the dot
                        _reset_for_new_input(picked, input_type)
            except Exception as e:
                st.error(str(e))

        input_path = st.session_state.get("input_path")
        input_type = st.session_state.get("input_type")
        if input_path:
            st.success(f"Selected: {input_path}")
            st.caption(f"Type: {input_type.upper()}")

        # PPTX conversion info (LibreOffice is now the only method)
        if input_type == "pptx":
            st.caption("📄 Using LibreOffice for PPTX conversion")
            libreoffice_timeout_seconds = int(
                st.number_input(
                    "LibreOffice timeout (seconds)",
                    min_value=30,
                    max_value=1800,
                    value=int(st.session_state.get("libreoffice_timeout_seconds", 120)),
                    step=30,
                    help="Increase this for large/complex decks that may exceed default conversion time.",
                )
            )
            st.session_state["libreoffice_timeout_seconds"] = libreoffice_timeout_seconds

        st.subheader("2) Process Input")
        # Render slides label
        pptx_label = "Render slides with LibreOffice" if input_type == "pptx" else "Process file"
        
        process_label = {
            "pptx": pptx_label,
            "pdf": "Extract PDF pages",
            "gif": "Extract GIF frames",
            "txt": "Convert text to images",
            "docx": "Extract DOCX text to images",
            "image": "Process image",
            "md": "Convert Markdown to images",
        }.get(input_type, "Process file")
        
        if st.button(process_label, use_container_width=True, disabled=not bool(input_path)):
            try:
                out_dir = Path(__file__).resolve().parent / "pptx2nano_output_streamlit"
                rendered_dir = out_dir / Path(input_path).stem / "rendered"
                
                if input_type == "pptx":
                    with st.spinner("Rendering slides using LibreOffice…"):
                        rendered_paths = export_slides(
                            Path(input_path),
                            rendered_dir,
                            method="libreoffice",
                            timeout_seconds=int(st.session_state.get("libreoffice_timeout_seconds", 120)),
                        )
                elif input_type == "pdf":
                    with st.spinner("Extracting PDF pages…"):
                        rendered_paths = _process_pdf_input(Path(input_path), rendered_dir)
                elif input_type == "gif":
                    with st.spinner("Extracting GIF frames…"):
                        rendered_paths = _process_gif_input(Path(input_path), rendered_dir)
                elif input_type == "txt":
                    with st.spinner("Converting text to images…"):
                        rendered_paths = _process_text_input(Path(input_path), rendered_dir)
                elif input_type == "docx":
                    with st.spinner("Extracting DOCX text…"):
                        rendered_paths = _process_docx_input(Path(input_path), rendered_dir)
                elif input_type == "image":
                    with st.spinner("Processing image…"):
                        rendered_paths = _process_image_input(Path(input_path), rendered_dir)
                elif input_type == "md":
                    with st.spinner("Converting Markdown…"):
                        rendered_paths = _process_markdown_input(Path(input_path), rendered_dir)
                else:
                    st.error(f"Unsupported file type: {input_type}")
                    return
                
                st.session_state["rendered_paths"] = rendered_paths
                if input_type == "pptx":
                    # Warm the text cache once so downstream generation avoids repeated parsing.
                    _get_pptx_text_cache(Path(input_path))
                st.success(f"Processed {len(rendered_paths)} pages/slides/frames.")
            except Exception as e:
                st.error(str(e))

        rendered_paths = st.session_state.get("rendered_paths")
        if rendered_paths:
            st.caption(f"Pages rendered: {len(rendered_paths)}")

    with left:
        st.subheader("Style Example")
        st.session_state["style_example_set"] = st.selectbox(
            "Example set",
            options=["slide1", "slide2"],
            index=0 if st.session_state.get("style_example_set", "slide1") == "slide1" else 1,
            help="Switch which cached style thumbnails you want to see.",
        )

        # Style selection setup
        style_options = _style_options()
        current_style = st.session_state.get("selected_style", "lego")
        
        style_choice = current_style
        custom_style = ""
        if style_choice == "custom":
            custom_style = st.text_input("Custom style", value="")
        style = _get_selected_style(style_choice, custom_style)

        image_model = pptx2nano.DEFAULT_IMAGE_MODEL

        if style_choice != "custom":
            st.caption(pptx2nano.BUILTIN_STYLES.get(style_choice, ""))

        # Display style example image
        if style:
            try:
                with st.spinner("Loading style example…"):
                    example_bytes = _get_or_create_style_example(
                        style,
                        image_model=image_model,
                        style_example_set=st.session_state.get("style_example_set", "slide1"),
                    )
                st.image(example_bytes, caption=f"Example: {style}", use_container_width=True)
            except Exception as e:
                st.warning(f"Could not load style example: {e}")

    with right:
        st.subheader("Style")
        st.caption("Choose a style:")
        cols = st.columns(2)
        
        for i, style_name in enumerate(style_options):
            with cols[i % 2]:
                # Use buttons instead of radio buttons for better control
                button_type = "primary" if style_name == current_style else "secondary"
                if st.button(style_name, key=f"style_btn_{i}", type=button_type, use_container_width=True):
                    st.session_state["selected_style"] = style_name
                    st.rerun()

    st.divider()

    input_path = st.session_state.get("input_path")
    rendered_paths = st.session_state.get("rendered_paths")
    if not input_path or not rendered_paths:
        st.info("Select a file and process it to begin.")
        return

    total_slides = len(rendered_paths)
    slide_numbers = list(range(1, total_slides + 1))

    st.subheader("3) Generate / Regenerate a slide")
    controls_col, action_col = st.columns([1, 2])

    with controls_col:
        # Get current slide index to maintain scroll position
        slide_options = ["ALL"] + slide_numbers
        current_slide = st.session_state.get("selected_slide", 1)
        try:
            if current_slide == "ALL":
                current_slide_index = 0
            else:
                current_slide_index = slide_options.index(current_slide)
        except ValueError:
            current_slide_index = 1
        
        gen_target = st.selectbox("Slide", slide_options, index=current_slide_index)
        st.session_state["selected_slide"] = gen_target
        
        overwrite_existing = st.checkbox(
            "Overwrite slides that were already generated",
            value=False,
            help="When unchecked, Generate ALL will skip slides that already have a generated image so you can keep different styles per slide.",
        )

        generate_workers = int(
            st.number_input(
                "Generate ALL workers",
                min_value=1,
                max_value=16,
                value=int(st.session_state.get("generate_workers", 4)),
                step=1,
                help="Parallel workers for Generate ALL. Higher can be faster but may hit API limits.",
            )
        )
        st.session_state["generate_workers"] = generate_workers

        text_extraction_mode = st.session_state.get("text_extraction_mode", "strict")
        dedupe_extracted_text = bool(st.session_state.get("dedupe_extracted_text", False))
        if input_type == "pptx":
            text_extraction_mode = st.selectbox(
                "Text extraction mode (PPTX)",
                options=["off", "strict", "assisted"],
                index=["off", "strict", "assisted"].index(
                    st.session_state.get("text_extraction_mode", "strict")
                ),
                help=(
                    "off: image-only. strict: use extracted text as-is for higher fidelity. "
                    "assisted: use Gemini cleanup for potentially better readability."
                ),
            )
            st.session_state["text_extraction_mode"] = text_extraction_mode
            dedupe_extracted_text = st.checkbox(
                "Deduplicate extracted text lines",
                value=bool(st.session_state.get("dedupe_extracted_text", False)),
                help="Turn on only if your source slides commonly duplicate identical lines unintentionally.",
            )
            st.session_state["dedupe_extracted_text"] = dedupe_extracted_text

    # When ALL is selected, we don't show a single rendered_path preview.
    slide_n = None if gen_target == "ALL" else int(gen_target)
    rendered_path = None if slide_n is None else rendered_paths[slide_n - 1]

    with action_col:
        pptx_input_path = Path(input_path) if input_type == "pptx" else None
        pptx_slide_texts: list[str] = []
        if pptx_input_path and text_extraction_mode != "off":
            try:
                pptx_slide_texts = _get_pptx_text_cache(pptx_input_path)
            except Exception:
                pptx_slide_texts = []

        if slide_n is None:
            if st.button("Generate ALL slides (skip existing unless overwrite checked)", use_container_width=True):
                try:
                    t0 = time.time()
                    durations: list[float] = []
                    prog = st.progress(0)
                    status = st.empty()

                    total = len(slide_numbers)
                    done = 0
                    to_generate: list[int] = []

                    for n in slide_numbers:
                        already = n in st.session_state["generated_images"]
                        if already and not overwrite_existing:
                            done += 1
                            prog.progress(int((done / total) * 100))
                            status.caption(f"Skipped existing slide {n} ({done}/{total})")
                            continue
                        to_generate.append(n)

                    if not to_generate:
                        status.success("Nothing to generate. All slides already existed and overwrite was off.")
                    else:
                        start_times: dict[int, float] = {}
                        errors: list[tuple[int, str]] = []

                        with concurrent.futures.ThreadPoolExecutor(
                            max_workers=max(1, min(generate_workers, len(to_generate)))
                        ) as ex:
                            future_to_slide = {}
                            for n in to_generate:
                                rp = rendered_paths[n - 1]
                                prev_id = st.session_state["interaction_ids"].get(n)
                                extracted_text = None
                                if pptx_slide_texts and (n - 1) < len(pptx_slide_texts):
                                    extracted_text = pptx_slide_texts[n - 1]
                                start_times[n] = time.time()
                                future = ex.submit(
                                    _call_slide_image_model,
                                    slide_index_1based=n,
                                    rendered_path=rp,
                                    image_model=image_model,
                                    style=style,
                                    total_slides=total_slides,
                                    previous_interaction_id=prev_id,
                                    pptx_path=pptx_input_path,
                                    text_extraction_mode=text_extraction_mode,
                                    dedupe_extracted_text=dedupe_extracted_text,
                                    extracted_slide_text=extracted_text,
                                )
                                future_to_slide[future] = n

                            for fut in concurrent.futures.as_completed(future_to_slide):
                                n = future_to_slide[fut]
                                try:
                                    img_bytes, new_id = fut.result()
                                except Exception as e:
                                    errors.append((n, str(e)))
                                else:
                                    st.session_state["generated_images"][n] = img_bytes
                                    st.session_state["interaction_ids"][n] = new_id
                                    st.session_state["last_preview_pdf"] = None
                                    dt = time.time() - start_times.get(n, time.time())
                                    durations.append(dt)

                                done += 1
                                avg = sum(durations) / len(durations) if durations else 0.0
                                remaining = max(0, total - done)
                                eta = remaining * avg
                                prog.progress(int((done / total) * 100))
                                if errors:
                                    status.caption(
                                        f"Processed slide {n} ({done}/{total}) | errors={len(errors)} "
                                        f"avg={avg:.1f}s ETA={eta:.0f}s"
                                    )
                                else:
                                    last = durations[-1] if durations else 0.0
                                    status.caption(
                                        f"Generated slide {n} ({done}/{total}) | "
                                        f"last={last:.1f}s avg={avg:.1f}s ETA={eta:.0f}s"
                                    )

                    prog.progress(100)
                    total_time = time.time() - t0
                    if "errors" in locals() and errors:
                        first_slide, first_error = errors[0]
                        status.error(
                            f"Generate ALL finished in {total_time:.1f}s with {len(errors)} errors. "
                            f"First failure: slide {first_slide} -> {first_error}"
                        )
                    else:
                        status.success(f"Done. Generate ALL finished in {total_time:.1f}s")
                except Exception as e:
                    st.error(str(e))
        else:
            prev_id = st.session_state["interaction_ids"].get(slide_n)
            if st.button("Generate / Regenerate this slide", use_container_width=True):
                try:
                    extracted_text = None
                    if pptx_slide_texts and (slide_n - 1) < len(pptx_slide_texts):
                        extracted_text = pptx_slide_texts[slide_n - 1]
                    prog = st.progress(0)
                    with st.spinner("Calling Gemini…"):
                        prog.progress(20)
                        img_bytes, new_id = _call_slide_image_model(
                            slide_index_1based=slide_n,
                            rendered_path=rendered_path,
                            image_model=image_model,
                            style=style,
                            total_slides=total_slides,
                            previous_interaction_id=prev_id,
                            pptx_path=pptx_input_path,
                            text_extraction_mode=text_extraction_mode,
                            dedupe_extracted_text=dedupe_extracted_text,
                            extracted_slide_text=extracted_text,
                        )
                        prog.progress(90)
                        st.session_state["generated_images"][slide_n] = img_bytes
                        st.session_state["interaction_ids"][slide_n] = new_id
                        st.session_state["last_preview_pdf"] = None
                        prog.progress(100)
                    st.success(f"Generated slide {slide_n}.")
                except Exception as e:
                    st.error(str(e))

    if slide_n is not None and rendered_path is not None:
        st.caption("Exhibit comparison")
        original_col, generated_col = st.columns([1, 1])
        with original_col:
            st.image(
                str(rendered_path),
                caption=f"Original slide {slide_n}",
                use_container_width=True,
            )
        with generated_col:
            gen_bytes = st.session_state["generated_images"].get(slide_n)
            if gen_bytes:
                st.image(
                    gen_bytes,
                    caption=f"Generated slide {slide_n}",
                    use_container_width=True,
                )
            else:
                st.info("No generated image for this slide yet.")

    st.divider()

    st.subheader("4) Preview and Save PDF")

    preview_max_pages = int(
        st.number_input(
            "Preview max pages",
            min_value=1,
            max_value=100,
            value=int(st.session_state.get("preview_max_pages", 10)),
            step=1,
            help="Only the first N pages are rendered for preview speed.",
        )
    )
    st.session_state["preview_max_pages"] = preview_max_pages

    preview_mode = st.selectbox(
        "Include slides",
        options=["ALL", "Custom"],
        index=0,
    )
    if preview_mode == "ALL":
        include_slides_sorted = slide_numbers
        st.caption("Preview will include all slides (that have generated images).")
    else:
        include_slides = st.multiselect(
            "Slides to include in PDF",
            options=slide_numbers,
            default=slide_numbers,
        )
        include_slides_sorted = sorted(include_slides)

    b1, b2 = st.columns([1, 1])

    with b1:
        if st.button("Build PDF preview", use_container_width=True):
            try:
                with st.spinner("Building PDF preview…"):
                    pdf_bytes = _build_pdf_bytes(include_slides_sorted)
                    st.session_state["last_preview_pdf"] = pdf_bytes
                st.success("Preview ready.")
            except Exception as e:
                st.error(str(e))

    with b2:
        preview_bytes = st.session_state.get("last_preview_pdf")
        can_save = bool(preview_bytes) and bool(input_path)
        if st.button("Save PDF next to input file (overwrite)", use_container_width=True, disabled=not can_save):
            try:
                input_p = Path(input_path)
                out_path = input_p.parent / f"{input_p.stem}_nano.pdf"
                out_path.write_bytes(preview_bytes)
                st.success(f"Saved: {out_path}")
            except Exception as e:
                st.error(str(e))

    preview_bytes = st.session_state.get("last_preview_pdf")
    if preview_bytes:
        _render_pdf_inline(preview_bytes, max_pages=preview_max_pages)
        st.download_button(
            label="Download preview PDF",
            data=preview_bytes,
            file_name=f"{Path(input_path).stem}_nano.pdf",
            mime="application/pdf",
            use_container_width=True,
        )


if __name__ == "__main__":
    main()
